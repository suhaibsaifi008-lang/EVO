"""Office & data intelligence — read/edit DOCX/XLSX/PPTX, analyze spreadsheets."""
import os
from datetime import datetime
from pathlib import Path

from .config import DATA_DIR

DOCS_DIR = DATA_DIR / "documents"
DOCS_DIR.mkdir(parents=True, exist_ok=True)


class DocError(RuntimeError):
    pass


def _resolve(path: str) -> Path:
    p = Path(path).expanduser()
    if not p.is_absolute():
        candidate = DATA_DIR.parent / "data" / path  # legacy relative
        p = candidate if candidate.exists() else Path.cwd() / path
    if not p.exists():
        raise DocError(f"file not found: {path}")
    return p


def read_office(path: str, max_chars: int = 4000) -> str:
    p = _resolve(path)
    ext = p.suffix.lower()
    if ext == ".docx":
        import docx

        d = docx.Document(str(p))
        parts = [par.text for par in d.paragraphs if par.text.strip()]
        for table in d.tables[:5]:
            for row in table.rows[:10]:
                cells = [c.text.strip() for c in row.cells]
                parts.append(" | ".join(cells))
        text = "\n".join(parts)
    elif ext in (".xlsx", ".xlsm"):
        import openpyxl

        parts = []
        with openpyxl.load_workbook(str(p), read_only=True, data_only=True) as wb:
            for ws in wb.worksheets[:6]:
                parts.append(f"[sheet: {ws.title}] {ws.max_row} rows x {ws.max_column} cols")
                for row in ws.iter_rows(min_row=1, max_row=8, values_only=True):
                    vals = [str(v) for v in row if v is not None][:12]
                    if vals:
                        parts.append("  " + " | ".join(vals))
        text = "\n".join(parts)
    elif ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            raise DocError("pip install python-pptx to read PowerPoint files")
        prs = Presentation(str(p))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        texts.append(t)
            if texts:
                parts.append(f"[slide {i}] " + " ".join(texts)[:400])
        text = "\n".join(parts)
    elif ext == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(p))
        text = "\n".join((page.extract_text() or "") for page in reader.pages[:40])
    else:
        text = p.read_text(encoding="utf-8", errors="ignore")
    return text[:max_chars]


def excel_profile(path: str) -> dict:
    import pandas as pd

    p = _resolve(path)
    xl = pd.ExcelFile(str(p))
    profile: dict = {"sheets": {}}
    for name in xl.sheet_names[:6]:
        df = xl.parse(name)
        info: dict = {"rows": int(len(df)), "columns": list(map(str, df.columns[:15]))}
        numerics = df.select_dtypes(include="number")
        if len(numerics.columns):
            col = numerics.columns[0]
            s = numerics[col].dropna()
            info["first_numeric_column"] = str(col)
            info["stats"] = {
                "count": int(s.count()),
                "mean": round(float(s.mean()), 2) if len(s) else None,
                "min": float(s.min()) if len(s) else None,
                "max": float(s.max()) if len(s) else None,
            }
            anomalies = int(len(s[(s - s.mean()).abs() > (3 * s.std())])) if s.std() and len(s) > 2 else 0
            info["possible_anomalies"] = anomalies
        profile["sheets"][name] = info
    return profile


def make_document(title: str, content: str, filename: str = "") -> str:
    import docx

    base = re_slug(filename or title)
    out = DOCS_DIR / f"{base}.docx"
    d = docx.Document()
    d.add_heading(title[:120], level=0)
    for block in (content or "").split("\n\n"):
        block = block.strip()
        if not block:
            continue
        if block.startswith("# "):
            d.add_heading(block[2:][:120], level=1)
            continue
        if block.startswith("## "):
            d.add_heading(block[3:][:120], level=2)
            continue
        lines = block.splitlines()
        bullets = [l[2:].strip() for l in lines if l.startswith(("- ", "* "))]
        if bullets and all(l.startswith(("- ", "* ")) for l in lines if l.strip()):
            for b in bullets:
                d.add_paragraph(b, style="List Bullet")
            continue
        d.add_paragraph(block.replace("\n", " "))
    d.save(str(out))
    return str(out)


def re_slug(name: str) -> str:
    import re

    base = re.sub(r"[^a-zA-Z0-9_-]+", "-", (name or "").strip())[:40].strip("-")
    return base or f"doc-{datetime.now().strftime('%Y%m%d_%H%M%S')}"
